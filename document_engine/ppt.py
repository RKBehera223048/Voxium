import os
import re
from typing import Optional
from pptx import Presentation
from .exporter import BaseExporter

class PptExporter(BaseExporter):
    """Exports markdown content to PowerPoint."""
    
    def export(self, content: str, output_path: str, title: Optional[str] = None) -> str:
        prs = Presentation()
        
        if title:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title
            
        lines = content.split('\n')
        current_slide = None
        current_tf = None
        
        for line in lines:
            heading_match = re.match(r'^(#{1,2})\s+(.*)$', line)
            if heading_match:
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                title_shape = shapes.title
                title_shape.text = heading_match.group(2)
                body_shape = shapes.placeholders[1]
                current_tf = body_shape.text_frame
                current_tf.text = ""
                continue
                
            if not current_slide:
                continue
                
            if line.startswith('- ') or line.startswith('* '):
                p = current_tf.add_paragraph()
                p.text = line[2:]
                p.level = 0
            elif line.startswith('  - '):
                p = current_tf.add_paragraph()
                p.text = line[4:]
                p.level = 1
            elif line.strip() != '' and not line.startswith('```'):
                p = current_tf.add_paragraph()
                p.text = line
                p.level = 0
                
        prs.save(output_path)
        return os.path.abspath(output_path)
